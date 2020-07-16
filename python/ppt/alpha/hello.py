# https://python-pptx.readthedocs.io/en/latest/user/quickstart.html

from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')

# l127914@pn1249300.lanl.gov:alpha $ py hello.py
# Traceback (most recent call last):
#   File "hello.py", line 3, in <module>
#     from pptx import Presentation
# ModuleNotFoundError: No module named 'pptx'

# l127914@pn1249300.lanl.gov:alpha $ date
# Wed Dec 19 15:21:35 MST 2018

# l127914@pn1249300.lanl.gov:alpha $ pwd
# /Volumes/Tlaltecuhtli/repos/GitHub/topa-development/python/ppt/alpha
